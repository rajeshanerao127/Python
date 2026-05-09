from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as ns
from lxml import etree

# Color Palette - Maroon & Gold University Style
MAROON = RGBColor(0x7B, 0x00, 0x20)       # Deep maroon
GOLD = RGBColor(0xC9, 0xA0, 0x2C)          # Rich gold
LIGHT_GOLD = RGBColor(0xF5, 0xE6, 0xB0)    # Light gold/cream
DARK_MAROON = RGBColor(0x55, 0x00, 0x14)   # Darker maroon
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x2D)
LIGHT_BG = RGBColor(0xFD, 0xF8, 0xF0)      # Warm white bg
MID_GOLD = RGBColor(0xE8, 0xCB, 0x6B)      # Medium gold accent

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

W = 13.33
H = 7.5

def hex_to_rgb(hex_str):
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, color, transparency=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, x, y, w, h, font_size=16, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False, font_name="Georgia", wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_bullet_text(slide, items, x, y, w, h, font_size=14, color=DARK_GRAY, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, ns.qn('a:buChar'))
        buChar.set('char', '◆')
        buSzPct = etree.SubElement(pPr, ns.qn('a:buSz'))
        # set bullet color
        buClr = etree.SubElement(pPr, ns.qn('a:buClr'))
        srgbClr = etree.SubElement(buClr, ns.qn('a:srgbClr'))
        srgbClr.set('val', 'C9A02C')
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = font_name
        # paragraph spacing
        pPr.set(ns.qn('a:spcBef'), '0')
        spcBef = etree.SubElement(pPr, ns.qn('a:spcBef'))
        spcPct = etree.SubElement(spcBef, ns.qn('a:spcPts'))
        spcPct.set('val', '80')
    return txBox

def slide_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_divider(slide, x, y, w):
    """Add a gold horizontal line"""
    line = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

# ─── SLIDE 1: TITLE SLIDE ───────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide, DARK_MAROON)

# Gold top band
add_rect(slide, 0, 0, W, 0.12, GOLD)
# Gold bottom band
add_rect(slide, 0, H-0.12, W, 0.12, GOLD)

# Left decorative block
add_rect(slide, 0, 0.12, 0.06, H-0.24, GOLD)

# Main maroon content area (slightly lighter)
add_rect(slide, 0.5, 0.8, 12.33, 5.9, RGBColor(0x8B, 0x10, 0x30))

# University name
add_text(slide, "SWAMI RAMANAND TEERTH MARATHWADA UNIVERSITY, NANDED",
         0.8, 0.95, 11.5, 0.6, font_size=14, bold=True, color=LIGHT_GOLD,
         align=PP_ALIGN.CENTER, font_name="Georgia")
add_text(slide, "SCHOOL OF PHARMACY",
         0.8, 1.5, 11.5, 0.4, font_size=13, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, font_name="Georgia")

add_divider(slide, 1.5, 2.0, 10.3)

# Main title
add_text(slide, "FORMULATION AND EVALUATION OF",
         0.8, 2.1, 11.5, 0.55, font_size=22, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, font_name="Georgia")
add_text(slide, "HERBAL COUGH RELIEF LOZENGES",
         0.8, 2.6, 11.5, 0.65, font_size=28, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Georgia")

add_text(slide, "Project Work (BP813PW) — B. Pharm Final Year (Sem VIII)",
         0.8, 3.25, 11.5, 0.45, font_size=13, bold=False, color=LIGHT_GOLD,
         align=PP_ALIGN.CENTER, font_name="Calibri")

add_divider(slide, 1.5, 3.8, 10.3)

# Student & Guide info
add_text(slide, "Submitted by:", 0.8, 3.95, 5.5, 0.35, font_size=11,
         bold=False, color=LIGHT_GOLD, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(slide, "Miss. Anerao Sakshi Mahendra [VZ15303]",
         0.8, 4.28, 5.5, 0.4, font_size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(slide, "B. Pharm, Fourth Year (Sem VIII)",
         0.8, 4.65, 5.5, 0.35, font_size=11, bold=False, color=LIGHT_GOLD,
         align=PP_ALIGN.CENTER, font_name="Calibri")

add_text(slide, "Under the Supervision of:", 6.8, 3.95, 5.5, 0.35, font_size=11,
         bold=False, color=LIGHT_GOLD, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(slide, "Dr. S. S. Pekamwar",
         6.8, 4.28, 5.5, 0.4, font_size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(slide, "Senior Professor, School of Pharmacy",
         6.8, 4.65, 5.5, 0.35, font_size=11, bold=False, color=LIGHT_GOLD,
         align=PP_ALIGN.CENTER, font_name="Calibri")

add_text(slide, "Academic Year: 2025–2026",
         0.8, 5.4, 11.5, 0.4, font_size=12, bold=True, color=GOLD,
         align=PP_ALIGN.CENTER, font_name="Calibri")

# ─── HELPER: content slide header ───────────────────────────────────────────
def content_header(slide, title, subtitle=None):
    slide_bg(slide, LIGHT_BG)
    add_rect(slide, 0, 0, W, 1.05, DARK_MAROON)
    add_rect(slide, 0, 1.05, W, 0.06, GOLD)
    add_text(slide, title, 0.4, 0.1, 12.5, 0.65, font_size=22, bold=True,
             color=WHITE, align=PP_ALIGN.LEFT, font_name="Georgia")
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.7, 12.5, 0.35, font_size=12,
                 color=LIGHT_GOLD, align=PP_ALIGN.LEFT, font_name="Calibri", italic=True)

def section_card(slide, x, y, w, h, title, body_lines, title_bg=MAROON):
    add_rect(slide, x, y, w, 0.42, title_bg)
    add_text(slide, title, x+0.12, y+0.05, w-0.24, 0.35,
             font_size=12, bold=True, color=WHITE, font_name="Georgia")
    add_rect(slide, x, y+0.42, w, h-0.42, WHITE)
    body_text = "\n".join(f"  ◆ {line}" for line in body_lines)
    tb = slide.shapes.add_textbox(Inches(x+0.12), Inches(y+0.48),
                                   Inches(w-0.24), Inches(h-0.55))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"◆  {line}"
        run.font.size = Pt(11)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Calibri"
    return

# ─── SLIDE 2: TABLE OF CONTENTS ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Table of Contents", "Formulation & Evaluation of Herbal Cough Relief Lozenges")

modules = [
    ("01", "Pharmaceutical Excipient Profile",
     ["Tulsi (Holy Basil)", "Liquorice (Mulethi)", "Ginger (Adrak)", "Mint (Pudina)"]),
    ("02", "Preformulation Studies of Tulsi",
     ["Biological Source & Classification", "Chemical Constituents", "Pharmacological Activities", "Dosage Forms"]),
    ("03", "Formulation & Evaluation of Herbal Lozenges",
     ["Aim, Types, Advantages", "Drug Selection & Composition", "Preparation Method", "Evaluation Parameters & Results"]),
]

for i, (num, title, sub) in enumerate(modules):
    yt = 1.45 + i * 1.85
    # Module number circle
    add_rect(slide, 0.5, yt, 1.0, 1.5, DARK_MAROON)
    add_text(slide, num, 0.5, yt+0.3, 1.0, 0.7, font_size=26, bold=True,
             color=GOLD, align=PP_ALIGN.CENTER, font_name="Georgia")
    # Title bar
    add_rect(slide, 1.6, yt, 11.2, 0.45, MAROON)
    add_text(slide, f"Module {num}: {title}", 1.75, yt+0.06, 11.0, 0.35,
             font_size=14, bold=True, color=WHITE, font_name="Georgia")
    # Sub items
    add_rect(slide, 1.6, yt+0.45, 11.2, 1.05, WHITE)
    cols = ["  ◆  " + sub[0] + "     ◆  " + sub[1],
            "  ◆  " + sub[2] + "     ◆  " + sub[3]]
    for j, col in enumerate(cols):
        add_text(slide, col, 1.75, yt+0.5+j*0.45, 11.0, 0.42,
                 font_size=11, color=DARK_GRAY, font_name="Calibri")

# ─── SLIDE 3: MODULE 1 INTRO ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Module 01: Pharmaceutical Excipient Profile", "Understanding the herbal ingredients used in lozenges")

add_text(slide,
    "Pharmaceutical excipients are substances other than the pharmacologically active drug included in the manufacturing "
    "process or contained in a finished pharmaceutical product. In this project, herbal excipients serve both as active "
    "therapeutic agents and as formulation components.",
    0.5, 1.3, 12.3, 1.0, font_size=13, color=DARK_GRAY, font_name="Calibri")

herbs = [
    ("TULSI", "Ocimum sanctum Linn.", "Lamiaceae", "Anti-tussive, Antioxidant,\nAntimicrobial, Immunity boost"),
    ("LIQUORICE", "Glycyrrhiza glabra Linn.", "Leguminacae", "Cough relief, Anti-inflammatory,\nFlavoring agent in lozenges"),
    ("GINGER", "Zingiber officinale", "Zingiberaceae", "Anti-nausea, Digestive aid,\nImmune health vs cold"),
    ("MINT", "Mentha spicata/piperita", "Lamiaceae", "Cooling agent, Flavoring,\nRespiratory relief"),
]

card_w = 2.9
for i, (name, bio, fam, uses) in enumerate(herbs):
    cx = 0.5 + i * 3.1
    cy = 2.55
    add_rect(slide, cx, cy, card_w, 0.55, DARK_MAROON)
    add_text(slide, name, cx+0.1, cy+0.08, card_w-0.2, 0.4,
             font_size=14, bold=True, color=GOLD, font_name="Georgia")
    add_rect(slide, cx, cy+0.55, card_w, 3.1, WHITE)
    # Shadow
    shadow_shape = slide.shapes.add_shape(1, Inches(cx+0.05), Inches(cy+0.65),
                                           Inches(card_w-0.1), Inches(3.0))
    shadow_shape.fill.background()
    shadow_shape.line.color.rgb = RGBColor(0xE0, 0xD0, 0xA0)
    shadow_shape.line.width = Pt(1)

    add_text(slide, "Biological Source:", cx+0.15, cy+0.62, card_w-0.3, 0.3,
             font_size=10, bold=True, color=MAROON, font_name="Calibri")
    add_text(slide, bio, cx+0.15, cy+0.9, card_w-0.3, 0.35,
             font_size=10, italic=True, color=DARK_GRAY, font_name="Calibri")
    add_text(slide, "Family:", cx+0.15, cy+1.22, card_w-0.3, 0.28,
             font_size=10, bold=True, color=MAROON, font_name="Calibri")
    add_text(slide, fam, cx+0.15, cy+1.48, card_w-0.3, 0.28,
             font_size=10, color=DARK_GRAY, font_name="Calibri")
    add_text(slide, "Key Uses:", cx+0.15, cy+1.8, card_w-0.3, 0.28,
             font_size=10, bold=True, color=MAROON, font_name="Calibri")
    add_text(slide, uses, cx+0.15, cy+2.06, card_w-0.3, 0.9,
             font_size=10, color=DARK_GRAY, font_name="Calibri")

# ─── SLIDE 4: Excipient Details - Chemical Constituents ──────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Herbal Excipients: Chemical Constituents", "Module 01 — Key bioactive compounds")

data = [
    ("TULSI\n(Holy Basil)", "Eugenol, Apigenin, Orientin, Oleanolic acid,\nCarvacrol, Alkaloids, Tannins, Saponins, Flavonoids"),
    ("LIQUORICE\n(Mulethi)", "Glycyrrhizin, Glycyrrhetinic acid, Triterpene,\nFlavonoids (Liquiritin, Isoliquiritin), Coumarins, Volatile oils"),
    ("GINGER\n(Adrak)", "Gingerol, Zingiberene (aromatic terpene),\nZingeerone, Paradols, Volatile oils"),
    ("MINT\n(Pudina)", "Menthol, Menthone, Limonene, Carvone,\nPhenolic acid, Flavonoids"),
]

for i, (herb, constituents) in enumerate(data):
    yt = 1.35 + i * 1.45
    add_rect(slide, 0.5, yt, 2.5, 1.2, MAROON)
    add_text(slide, herb, 0.6, yt+0.2, 2.3, 0.85,
             font_size=12, bold=True, color=WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
    add_rect(slide, 3.1, yt, 9.7, 1.2, WHITE)
    # border
    border = slide.shapes.add_shape(1, Inches(3.1), Inches(yt), Inches(9.7), Inches(1.2))
    border.fill.background()
    border.line.color.rgb = RGBColor(0xD4, 0xB8, 0x60)
    border.line.width = Pt(1.5)
    add_text(slide, constituents, 3.25, yt+0.2, 9.4, 0.85,
             font_size=12, color=DARK_GRAY, font_name="Calibri")

# footer note
add_rect(slide, 0.5, 6.95, 12.3, 0.38, DARK_MAROON)
add_text(slide, "These herbal constituents collectively provide anti-tussive, anti-inflammatory, antimicrobial & soothing properties in lozenges.",
         0.7, 6.98, 12.0, 0.32, font_size=10, color=LIGHT_GOLD, font_name="Calibri")

# ─── SLIDE 5: MODULE 2 - Tulsi Overview ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Module 02: Preformulation Studies of Tulsi", "Ocimum sanctum Linn. — Holy Basil")

# Left column: Classification
add_rect(slide, 0.5, 1.28, 5.8, 5.85, WHITE)
border = slide.shapes.add_shape(1, Inches(0.5), Inches(1.28), Inches(5.8), Inches(5.85))
border.fill.background()
border.line.color.rgb = GOLD
border.line.width = Pt(1.5)

add_rect(slide, 0.5, 1.28, 5.8, 0.48, MAROON)
add_text(slide, "Scientific Classification", 0.65, 1.33, 5.5, 0.4,
         font_size=13, bold=True, color=WHITE, font_name="Georgia")

classification = [
    ("Kingdom", "Plantae"),
    ("Division", "Magnoliophyta"),
    ("Class", "Magnoliopsida"),
    ("Order", "Lamiales"),
    ("Family", "Lamiaceae"),
    ("Genus", "Ocimum"),
    ("Species", "Ocimum sanctum Linn."),
]
for j, (k, v) in enumerate(classification):
    ypos = 1.95 + j * 0.7
    bg_col = LIGHT_BG if j % 2 == 0 else WHITE
    add_rect(slide, 0.5, ypos-0.1, 5.8, 0.68, bg_col)
    add_text(slide, k + ":", 0.65, ypos, 2.0, 0.4, font_size=12,
             bold=True, color=MAROON, font_name="Calibri")
    add_text(slide, v, 2.8, ypos, 3.2, 0.4, font_size=12,
             color=DARK_GRAY, font_name="Calibri", italic=(k=="Species"))

# Right column: About Tulsi
add_rect(slide, 6.8, 1.28, 6.0, 5.85, WHITE)
border2 = slide.shapes.add_shape(1, Inches(6.8), Inches(1.28), Inches(6.0), Inches(5.85))
border2.fill.background()
border2.line.color.rgb = GOLD
border2.line.width = Pt(1.5)

add_rect(slide, 6.8, 1.28, 6.0, 0.48, DARK_MAROON)
add_text(slide, "About Tulsi (Holy Basil)", 6.95, 1.33, 5.7, 0.4,
         font_size=13, bold=True, color=WHITE, font_name="Georgia")

about = (
    "Tulsi (Ocimum sanctum), also known as Holy Basil, is a sacred aromatic herb "
    "deeply rooted in Indian culture and Ayurveda. It is revered for its potent "
    "medicinal properties and spiritual significance."
)
add_text(slide, about, 6.95, 1.9, 5.7, 1.0, font_size=11.5, color=DARK_GRAY, font_name="Calibri")

add_text(slide, "Key Properties:", 6.95, 3.05, 5.7, 0.35,
         font_size=12, bold=True, color=MAROON, font_name="Calibri")

key_props = [
    "Adaptogen — helps body cope with stress",
    "Expectorant — clears mucus from airways",
    "Antimicrobial — fights bacteria & viruses",
    "Anti-inflammatory & Antioxidant",
    "Supports respiratory health (cough/cold)",
    "Blood purifier & immune booster",
]
for j, prop in enumerate(key_props):
    add_text(slide, f"◆  {prop}", 6.95, 3.45 + j*0.5, 5.7, 0.42,
             font_size=11, color=DARK_GRAY, font_name="Calibri")

# ─── SLIDE 6: Tulsi Chemical Constituents Table ───────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Tulsi: Chemical Constituents & Properties", "Module 02 — Preformulation Studies")

# Table data
table_data = [
    ["Compound", "Category", "Medicinal Properties"],
    ["Eugenol", "Phenolic Compound", "Antioxidant, Anti-inflammatory, Antiseptic"],
    ["Ursolic Acid", "Triterpenoid", "Antimicrobial, Anticancer, Anti-inflammatory"],
    ["Rosmarinic Acid", "Polyphenol", "Neuroprotective, Anti-allergic"],
    ["Linalool", "Terpene Alcohol", "Anxiolytic, Mild Sedative, Anti-inflammatory"],
    ["Carvacrol", "Monoterpenoid", "Antioxidant, Antimicrobial"],
    ["Apigenin", "Flavonoid", "Antioxidant, Anticancer, Anti-anxiety"],
    ["Orientin", "Flavonoid Glycoside", "Antioxidant, Radioprotective"],
    ["Vitamin C", "Vitamin", "Immune Support, Antioxidant"],
]

col_widths = [2.3, 2.3, 7.1]
row_h = 0.54

for row_i, row in enumerate(table_data):
    ypos = 1.28 + row_i * row_h
    for col_i, cell in enumerate(row):
        cx = 0.5 + sum(col_widths[:col_i])
        cw = col_widths[col_i]
        if row_i == 0:
            bg = DARK_MAROON
            fc = WHITE
            fsz = 12
            fb = True
        elif row_i % 2 == 1:
            bg = RGBColor(0xFD, 0xF2, 0xE0)
            fc = DARK_GRAY
            fsz = 11
            fb = False
        else:
            bg = WHITE
            fc = DARK_GRAY
            fsz = 11
            fb = False
        add_rect(slide, cx, ypos, cw, row_h, bg)
        b = slide.shapes.add_shape(1, Inches(cx), Inches(ypos), Inches(cw), Inches(row_h))
        b.fill.background()
        b.line.color.rgb = RGBColor(0xD4, 0xB0, 0x50)
        b.line.width = Pt(0.75)
        add_text(slide, cell, cx+0.1, ypos+0.08, cw-0.2, row_h-0.12,
                 font_size=fsz, bold=fb, color=fc if row_i != 0 else WHITE,
                 font_name="Calibri" if row_i != 0 else "Georgia")

# ─── SLIDE 7: Tulsi Pharmacological Activities ───────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Tulsi: Pharmacological Activities", "Module 02 — Preformulation Studies")

activities = [
    ("Antioxidant &\nAdaptogenic", "Protects cells from oxidative damage; helps body adapt to stress by balancing neuroendocrine & immune systems"),
    ("Anti-inflammatory\n& Analgesic", "Reduces inflammation and pain through compounds like Eugenol and Ursolic acid"),
    ("Antimicrobial", "Effective against bacteria, viruses, fungi, and protozoa; used for infections & as natural preservative"),
    ("Cardioprotective", "Lowers cholesterol & triglycerides, supports heart health, may prevent platelet aggregation"),
    ("Respiratory\nSupport", "Treats cough, cold, asthma, and bronchitis — especially when combined with honey & ginger"),
    ("Neuroprotective", "Enhances memory, sharpens senses, protects nerves; used for fatigue and insomnia"),
]

cols = 3
cw = 3.9
ch = 2.0
for i, (title, desc) in enumerate(activities):
    row = i // cols
    col = i % cols
    cx = 0.5 + col * (cw + 0.3)
    cy = 1.35 + row * (ch + 0.25)
    add_rect(slide, cx, cy, cw, 0.55, MAROON)
    add_text(slide, title, cx+0.12, cy+0.04, cw-0.24, 0.48,
             font_size=11, bold=True, color=GOLD, font_name="Georgia", align=PP_ALIGN.CENTER)
    add_rect(slide, cx, cy+0.55, cw, ch-0.55, WHITE)
    brd = slide.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(cw), Inches(ch))
    brd.fill.background()
    brd.line.color.rgb = GOLD
    brd.line.width = Pt(1.2)
    add_text(slide, desc, cx+0.12, cy+0.62, cw-0.24, ch-0.7,
             font_size=11, color=DARK_GRAY, font_name="Calibri")

# ─── SLIDE 8: MODULE 3 INTRO ──────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Module 03: Formulation of Herbal Cough Relief Lozenges", "Introduction & Background")

# Left panel
add_rect(slide, 0.5, 1.28, 7.5, 5.85, WHITE)
brd = slide.shapes.add_shape(1, Inches(0.5), Inches(1.28), Inches(7.5), Inches(5.85))
brd.fill.background()
brd.line.color.rgb = GOLD
brd.line.width = Pt(1.5)

add_rect(slide, 0.5, 1.28, 7.5, 0.48, MAROON)
add_text(slide, "What are Lozenges?", 0.65, 1.33, 7.2, 0.4,
         font_size=14, bold=True, color=WHITE, font_name="Georgia")

loz_intro = (
    "Lozenges are palatable solid unit dosage forms administered in the oral cavity. "
    "They are solid, medicated, flavored, and sweetened base dosage forms intended to be "
    "sucked and held in the mouth or pharynx.\n\n"
    "Benefits include: increased retention time in oral cavity, improved bioavailability, "
    "reduced gastric irritation, and bypassing of first-pass metabolism."
)
add_text(slide, loz_intro, 0.65, 1.9, 7.2, 1.8, font_size=12, color=DARK_GRAY, font_name="Calibri")

add_text(slide, "Types of Lozenges:", 0.65, 3.8, 7.2, 0.38,
         font_size=13, bold=True, color=MAROON, font_name="Georgia")
types = [
    "Hard Lozenges — candy base, long dissolution time",
    "Soft Lozenges — glycerin/PEG base, softer texture",
    "Chewable Lozenges — compressed with chewable base",
    "Local Action — Antiseptics, Decongestants",
    "Systemic Action — Vitamins, Nicotine replacement",
]
for j, t in enumerate(types):
    add_text(slide, f"◆  {t}", 0.65, 4.22 + j*0.48, 7.2, 0.42,
             font_size=11, color=DARK_GRAY, font_name="Calibri")

# Right panel
add_rect(slide, 8.4, 1.28, 4.6, 5.85, WHITE)
brd2 = slide.shapes.add_shape(1, Inches(8.4), Inches(1.28), Inches(4.6), Inches(5.85))
brd2.fill.background()
brd2.line.color.rgb = GOLD
brd2.line.width = Pt(1.5)
add_rect(slide, 8.4, 1.28, 4.6, 0.48, DARK_MAROON)
add_text(slide, "Herbal Advantage", 8.55, 1.33, 4.3, 0.4,
         font_size=14, bold=True, color=GOLD, font_name="Georgia")

advantages = [
    "Palatability — pleasant taste",
    "Ease of Administration",
    "Accurate Dosage",
    "Improved Patient Compliance",
    "Rapid Dissolution",
    "Reduced GI Irritation",
    "Innovative Drug Delivery",
    "Natural & Safe Ingredients",
    "Anti-inflammatory action",
    "Suitable for pediatric patients",
]
for j, adv in enumerate(advantages):
    add_text(slide, f"◆  {adv}", 8.55, 1.9 + j*0.5, 4.3, 0.42,
             font_size=11, color=DARK_GRAY, font_name="Calibri")

# ─── SLIDE 9: Drug Selection ──────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Drug Selection: Tulsi as Antitussive Agent", "Module 03 — Scientific rationale")

# Rationale
add_rect(slide, 0.5, 1.28, 12.3, 1.4, RGBColor(0xFF, 0xF5, 0xDC))
brd = slide.shapes.add_shape(1, Inches(0.5), Inches(1.28), Inches(12.3), Inches(1.4))
brd.fill.background()
brd.line.color.rgb = GOLD
brd.line.width = Pt(1.5)
add_text(slide, "Why Tulsi?", 0.65, 1.32, 3.0, 0.38, font_size=13, bold=True,
         color=MAROON, font_name="Georgia")
rationale = ("Eugenol — found at 40–71% concentration in Tulsi essential oil — reduces the urge to cough via central action, "
             "mediated by opioid and GABA-ergic systems. Other constituents such as Cineole, Thymol, and Ursolic acid "
             "further contribute to antitussive and anti-inflammatory effects.")
add_text(slide, rationale, 0.65, 1.72, 12.0, 0.9, font_size=12, color=DARK_GRAY, font_name="Calibri")

# Distribution table
add_rect(slide, 0.5, 2.85, 12.3, 0.48, DARK_MAROON)
for ci, hdr in enumerate(["Herb", "Key Active Compound", "Therapeutic Role in Lozenges"]):
    cx = 0.5 + ci * 4.1
    add_text(slide, hdr, cx+0.1, 2.88, 3.9, 0.4, font_size=12, bold=True,
             color=WHITE, font_name="Georgia")

rows = [
    ("Tulsi (Ocimum sanctum)", "Eugenol, Cineole, Thymol", "Primary antitussive agent"),
    ("Liquorice (Glycyrrhiza glabra)", "Glycyrrhizin, Flavonoids", "Demulcent & soothing action"),
    ("Ginger (Zingiber officinale)", "Gingerol, Zingiberene", "Anti-inflammatory support"),
    ("Mint (Mentha spp.)", "Menthol, Menthone", "Cooling & flavoring agent"),
]
for ri, (h, c, r) in enumerate(rows):
    ypos = 3.33 + ri * 0.65
    bg = RGBColor(0xFD, 0xF2, 0xE0) if ri % 2 == 0 else WHITE
    for ci, txt in enumerate([h, c, r]):
        cx = 0.5 + ci * 4.1
        add_rect(slide, cx, ypos, 4.1, 0.62, bg)
        brd = slide.shapes.add_shape(1, Inches(cx), Inches(ypos), Inches(4.1), Inches(0.62))
        brd.fill.background()
        brd.line.color.rgb = RGBColor(0xD4, 0xB0, 0x50)
        brd.line.width = Pt(0.75)
        add_text(slide, txt, cx+0.1, ypos+0.1, 3.9, 0.45, font_size=11, color=DARK_GRAY, font_name="Calibri")

# Bottom note
add_rect(slide, 0.5, 5.95, 12.3, 0.48, MAROON)
add_text(slide, "Tulsi acts as a natural antitussive. Native to the Indian subcontinent — found in India, Nepal, Sri Lanka & Malaysia.",
         0.65, 5.99, 12.0, 0.38, font_size=11, color=WHITE, font_name="Calibri")

# ─── SLIDE 10: Composition Table ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Lozenge Composition & Preparation Method", "Module 03 — Formulation Details")

# Composition table
add_rect(slide, 0.5, 1.28, 7.6, 0.48, DARK_MAROON)
for ci, hdr in enumerate(["Sr.", "Ingredient", "Quantity", "Role"]):
    widths = [0.5, 2.2, 1.5, 3.2]
    cx = 0.5 + sum(widths[:ci])
    add_text(slide, hdr, cx+0.08, 1.3, widths[ci]-0.1, 0.42,
             font_size=12, bold=True, color=WHITE, font_name="Georgia")

composition = [
    ("1", "Tulsi Extract", "25 ml", "Anti-tussive, Antimicrobial"),
    ("2", "Ginger Extract", "20 ml", "Anti-inflammatory"),
    ("3", "Liquorice Extract", "25 ml", "Demulcent, Soothing"),
    ("4", "Water", "50 ml", "Solvent"),
    ("5", "Sugar (Sucrose)", "90 g", "Base — provides hardness"),
    ("6", "Menthol", "2–3 ml", "Cooling effect"),
    ("7", "Citric Acid", "5 ml", "Flavour enhancer"),
    ("8", "Flavour (Lemon)", "2–3 ml", "Taste improvement"),
    ("9", "Colour", "q.s.", "Appearance"),
]
widths = [0.5, 2.2, 1.5, 3.2]
for ri, row in enumerate(composition):
    ypos = 1.76 + ri * 0.55
    bg = RGBColor(0xFD, 0xF2, 0xE0) if ri % 2 == 0 else WHITE
    for ci, cell in enumerate(row):
        cx = 0.5 + sum(widths[:ci])
        add_rect(slide, cx, ypos, widths[ci], 0.52, bg)
        brd = slide.shapes.add_shape(1, Inches(cx), Inches(ypos), Inches(widths[ci]), Inches(0.52))
        brd.fill.background()
        brd.line.color.rgb = RGBColor(0xD4, 0xB0, 0x50)
        brd.line.width = Pt(0.75)
        add_text(slide, cell, cx+0.08, ypos+0.08, widths[ci]-0.12, 0.38,
                 font_size=11, color=DARK_GRAY, font_name="Calibri")

# Preparation method - right side
add_rect(slide, 8.5, 1.28, 4.6, 6.1, WHITE)
brd = slide.shapes.add_shape(1, Inches(8.5), Inches(1.28), Inches(4.6), Inches(6.1))
brd.fill.background()
brd.line.color.rgb = GOLD
brd.line.width = Pt(1.5)
add_rect(slide, 8.5, 1.28, 4.6, 0.48, MAROON)
add_text(slide, "Preparation Method", 8.65, 1.33, 4.3, 0.4,
         font_size=13, bold=True, color=WHITE, font_name="Georgia")

steps = [
    "Weigh all ingredients accurately",
    "Heat sugar + water until candy base forms",
    "Cool mixture to 80–90°C",
    "Add herbal extracts + excipients",
    "Mix the mixture homogeneously",
    "Pour into molds for shaping",
    "Allow cooling and solidification",
    "Remove lozenges & pack in butter paper or aluminium foil",
]
for j, step in enumerate(steps):
    cy2 = 1.9 + j * 0.62
    add_rect(slide, 8.6, cy2+0.04, 0.4, 0.4, MAROON)
    add_text(slide, str(j+1), 8.6, cy2+0.06, 0.4, 0.32,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text(slide, step, 9.1, cy2+0.04, 3.85, 0.45,
             font_size=10.5, color=DARK_GRAY, font_name="Calibri")

# ─── SLIDE 11: Evaluation Parameters ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Evaluation Parameters", "Module 03 — Quality Testing of Herbal Lozenges")

params = [
    ("1. Physical\nAppearance", "Organoleptic properties:\nColor: Dark Brown\nTexture: Hard\nOdour: Honey\nShape: Hexagon\nTaste: Sweet & Bitter"),
    ("2. Hardness\nTest", "10 lozenges tested using Monsanto Hardness Tester\n\nAverage Hardness:\n17.35 kg/cm\n\n(Range: 16.09 – 18.07 kg/cm)"),
    ("3. Weight\nVariation", "20 lozenges tested\n\nAverage Weight:\n6.47 g\n\n(Total: 129.54 g)"),
    ("4. Thickness\nTest", "5 lozenges tested using Vernier Calliper\n\nAverage Thickness:\n11.64 mm\n\n(Range: 11.00 – 12.41 mm)"),
]

card_w = 2.95
for i, (param, result) in enumerate(params):
    cx = 0.5 + i * 3.1
    cy = 1.38

    add_rect(slide, cx, cy, card_w, 0.72, DARK_MAROON)
    add_text(slide, param, cx+0.12, cy+0.06, card_w-0.24, 0.62,
             font_size=12, bold=True, color=GOLD, font_name="Georgia", align=PP_ALIGN.CENTER)
    add_rect(slide, cx, cy+0.72, card_w, 4.9, WHITE)
    brd = slide.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(card_w), Inches(5.62))
    brd.fill.background()
    brd.line.color.rgb = GOLD
    brd.line.width = Pt(1.5)
    add_text(slide, result, cx+0.15, cy+0.85, card_w-0.3, 4.6,
             font_size=12, color=DARK_GRAY, font_name="Calibri")

# ─── SLIDE 12: Results - Hardness & Weight Chart ──────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Evaluation Results: Hardness & Weight Variation", "Module 03 — Quantitative Analysis")

# Hardness bar chart
hardness_vals = [17.01, 17.05, 18.07, 18.01, 17.08, 16.09, 18.07, 17.07, 18.04, 17.06]
chart_data_h = [{"name": "Hardness (kg/cm)", "labels": [f"L{i+1}" for i in range(10)],
                  "values": hardness_vals}]
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

cd_h = ChartData()
cd_h.categories = [f"L{i+1}" for i in range(10)]
cd_h.add_series("Hardness (kg/cm)", hardness_vals)

chart_h = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.35), Inches(5.9), Inches(4.3), cd_h
).chart
chart_h.has_title = True
chart_h.chart_title.text_frame.text = "Hardness Test Results"
chart_h.series[0].format.fill.solid()
chart_h.series[0].format.fill.fore_color.rgb = MAROON

# Weight variation chart
weights = [6.93, 6.97, 6.61, 6.82, 6.31, 6.40, 6.54, 6.98, 6.89, 6.04,
           6.78, 7.01, 6.18, 6.02, 6.82, 6.68, 6.64, 5.84, 6.70, 5.98]
cd_w = ChartData()
cd_w.categories = [f"L{i+1}" for i in range(20)]
cd_w.add_series("Weight (g)", weights)

chart_w = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(6.9), Inches(1.35), Inches(5.9), Inches(4.3), cd_w
).chart
chart_w.has_title = True
chart_w.chart_title.text_frame.text = "Weight Variation (20 Lozenges)"
chart_w.series[0].format.line.color.rgb = GOLD

# Summary stats
add_rect(slide, 0.5, 5.75, 5.9, 0.6, DARK_MAROON)
add_text(slide, "Avg Hardness: 17.35 kg/cm  |  Range: 16.09 – 18.07 kg/cm",
         0.65, 5.82, 5.6, 0.42, font_size=11, bold=True, color=WHITE, font_name="Calibri")
add_rect(slide, 6.9, 5.75, 5.9, 0.6, DARK_MAROON)
add_text(slide, "Avg Weight: 6.47 g  |  Variation range: –9.74% to +8.35%",
         7.05, 5.82, 5.6, 0.42, font_size=11, bold=True, color=WHITE, font_name="Calibri")

# ─── SLIDE 13: Thickness Results ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_header(slide, "Evaluation Results: Thickness Test", "Module 03 — Physical Characterization")

thickness_vals = [11.46, 12.23, 12.41, 11.12, 11.00]
cd_t = ChartData()
cd_t.categories = [f"Lozenge {i+1}" for i in range(5)]
cd_t.add_series("Thickness (mm)", thickness_vals)

chart_t = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.35), Inches(6.5), Inches(4.6), cd_t
).chart
chart_t.has_title = True
chart_t.chart_title.text_frame.text = "Thickness of 5 Lozenges (mm)"
chart_t.series[0].format.fill.solid()
chart_t.series[0].format.fill.fore_color.rgb = GOLD

# Results table on right
add_rect(slide, 7.5, 1.35, 5.3, 0.48, DARK_MAROON)
for ci, hdr in enumerate(["Lozenge", "Thickness (mm)"]):
    cx = 7.5 + ci * 2.65
    add_text(slide, hdr, cx+0.1, 1.38, 2.5, 0.4, font_size=12, bold=True,
             color=WHITE, font_name="Georgia")

thickness_data = [(f"Lozenge {i+1}", f"{v:.2f}") for i, v in enumerate(thickness_vals)]
for ri, (loz, thick) in enumerate(thickness_data):
    ypos = 1.83 + ri * 0.65
    bg = RGBColor(0xFD, 0xF2, 0xE0) if ri % 2 == 0 else WHITE
    for ci, cell in enumerate([loz, thick]):
        cx = 7.5 + ci * 2.65
        add_rect(slide, cx, ypos, 2.65, 0.62, bg)
        brd = slide.shapes.add_shape(1, Inches(cx), Inches(ypos), Inches(2.65), Inches(0.62))
        brd.fill.background()
        brd.line.color.rgb = RGBColor(0xD4, 0xB0, 0x50)
        brd.line.width = Pt(0.75)
        add_text(slide, cell, cx+0.1, ypos+0.1, 2.5, 0.42,
                 font_size=12, color=DARK_GRAY, font_name="Calibri")

# Average
add_rect(slide, 7.5, 5.1, 5.3, 0.58, MAROON)
add_text(slide, "Average Thickness = 11.64 mm", 7.65, 5.18, 5.0, 0.42,
         font_size=13, bold=True, color=WHITE, font_name="Georgia")

add_rect(slide, 0.5, 6.1, 6.5, 0.5, DARK_MAROON)
add_text(slide, "Average Thickness of 5 Lozenges = 58.22 / 5 = 11.64 mm",
         0.65, 6.16, 6.2, 0.38, font_size=11, bold=True, color=WHITE, font_name="Calibri")

# ─── SLIDE 14: Conclusion ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide, DARK_MAROON)
add_rect(slide, 0, 0, W, 0.1, GOLD)
add_rect(slide, 0, H-0.1, W, 0.1, GOLD)

add_text(slide, "CONCLUSION", 0.5, 0.25, 12.3, 0.7, font_size=28, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER, font_name="Georgia")
add_divider(slide, 1.5, 1.0, 10.3)

conclusions = [
    ("Feasible Formulation", "The herbal lozenge preparation is simple, time-saving, and organoleptically acceptable, especially for pediatric patients."),
    ("Natural & Effective", "Herbal ingredients (Tulsi, Liquorice, Ginger, Mint) provide anti-tussive, anti-inflammatory, and antioxidant properties as safe alternatives to conventional cough drops."),
    ("Ideal Dosage Form", "Lozenges offer enhanced bioavailability, patient compliance, reduced GI irritation, and immediate onset of action at low doses."),
    ("Evaluation Success", "Physical evaluation showed satisfactory hardness (17.35 kg/cm), weight (6.47 g) and thickness (11.64 mm) — meeting acceptable pharmaceutical standards."),
    ("Future Potential", "Herbal lozenges represent a viable, all-natural approach to respiratory health. Further research on dosage optimization and clinical comparison is recommended."),
]

for i, (title, body) in enumerate(conclusions):
    cy = 1.15 + i * 1.12
    add_rect(slide, 0.5, cy, 2.5, 0.95, RGBColor(0x8B, 0x10, 0x30))
    add_text(slide, title, 0.6, cy+0.1, 2.3, 0.78, font_size=11, bold=True,
             color=GOLD, font_name="Georgia", align=PP_ALIGN.CENTER)
    add_rect(slide, 3.1, cy, 9.7, 0.95, RGBColor(0x6E, 0x00, 0x1A))
    add_text(slide, body, 3.25, cy+0.1, 9.4, 0.78, font_size=11.5,
             color=WHITE, font_name="Calibri")

# ─── SLIDE 15: THANK YOU ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide, DARK_MAROON)
add_rect(slide, 0, 0, W, 0.1, GOLD)
add_rect(slide, 0, H-0.1, W, 0.1, GOLD)

# Center decoration
add_rect(slide, 3.5, 1.5, 6.3, 4.5, RGBColor(0x8B, 0x10, 0x30))
brd = slide.shapes.add_shape(1, Inches(3.5), Inches(1.5), Inches(6.3), Inches(4.5))
brd.fill.background()
brd.line.color.rgb = GOLD
brd.line.width = Pt(2.5)

add_text(slide, "THANK YOU", 3.6, 2.1, 6.1, 1.0, font_size=40, bold=True,
         color=GOLD, align=PP_ALIGN.CENTER, font_name="Georgia")
add_divider(slide, 4.5, 3.25, 4.3)
add_text(slide, "Miss. Anerao Sakshi Mahendra", 3.6, 3.4, 6.1, 0.5,
         font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(slide, "B. Pharm IV Year | Sem VIII | [VZ15303]", 3.6, 3.9, 6.1, 0.4,
         font_size=12, color=LIGHT_GOLD, align=PP_ALIGN.CENTER, font_name="Calibri")

add_text(slide, "Under the Guidance of Dr. S. S. Pekamwar", 3.6, 4.4, 6.1, 0.4,
         font_size=11.5, italic=True, color=LIGHT_GOLD, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(slide, "SRTMU, Nanded — School of Pharmacy | 2025–26",
         3.6, 4.85, 6.1, 0.4, font_size=11, color=GOLD,
         align=PP_ALIGN.CENTER, font_name="Calibri")

# Save
output_path = "/mnt/user-data/outputs/Herbal_Cough_Lozenges_PPT.pptx"
prs.save(output_path)
print("Saved:", output_path)